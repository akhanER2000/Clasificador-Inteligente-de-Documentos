"""
Extractor de texto desde documentos no-imagen (PDF, Word, Excel, TXT, CSV, PPT).
Convierte el contenido a texto plano para luego procesarlo con Groq.
"""
import os
import csv
import io


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """
    Extrae texto plano de un archivo según su extensión.
    
    Args:
        file_bytes: Contenido del archivo en bytes.
        filename: Nombre del archivo con extensión.
        
    Returns:
        Texto extraído del documento.
    """
    ext = filename.lower().rsplit(".", 1)[-1]
    
    extractors = {
        "pdf": _extract_from_pdf,
        "docx": _extract_from_docx,
        "doc": _extract_from_docx,
        "xlsx": _extract_from_xlsx,
        "xls": _extract_from_xlsx,
        "pptx": _extract_from_pptx,
        "ppt": _extract_from_pptx,
        "txt": _extract_from_txt,
        "csv": _extract_from_csv,
    }
    
    extractor = extractors.get(ext)
    if not extractor:
        return f"[Formato .{ext} no soportado para extracción de texto]"
    
    try:
        return extractor(file_bytes)
    except Exception as e:
        return f"[Error al extraer texto: {str(e)}]"


def _extract_from_pdf(file_bytes: bytes) -> str:
    """Extrae texto de un PDF."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n\n".join(text_parts) if text_parts else "[PDF sin texto extraíble - puede ser un PDF escaneado]"
    except ImportError:
        return "[PyPDF2 no instalado - no se puede leer PDF]"


def _extract_from_docx(file_bytes: bytes) -> str:
    """Extrae texto de un archivo Word (.docx)."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        # También extraer tablas
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        return "\n".join(text_parts) if text_parts else "[Documento Word vacío]"
    except ImportError:
        return "[python-docx no instalado - no se puede leer Word]"


def _extract_from_xlsx(file_bytes: bytes) -> str:
    """Extrae texto de un archivo Excel (.xlsx)."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"=== Hoja: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    text_parts.append(" | ".join(row_values))
        wb.close()
        return "\n".join(text_parts) if text_parts else "[Excel vacío]"
    except ImportError:
        return "[openpyxl no instalado - no se puede leer Excel]"


def _extract_from_pptx(file_bytes: bytes) -> str:
    """Extrae texto de un archivo PowerPoint (.pptx)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text)
            if slide_texts:
                text_parts.append(f"--- Diapositiva {i} ---")
                text_parts.extend(slide_texts)
        return "\n".join(text_parts) if text_parts else "[Presentación vacía]"
    except ImportError:
        return "[python-pptx no instalado - no se puede leer PowerPoint]"


def _extract_from_txt(file_bytes: bytes) -> str:
    """Extrae texto de un archivo de texto plano."""
    # Intentar diferentes encodings
    for encoding in ["utf-8", "latin-1", "cp1252"]:
        try:
            return file_bytes.decode(encoding)
        except (UnicodeDecodeError, AttributeError):
            continue
    return file_bytes.decode("utf-8", errors="replace")


def _extract_from_csv(file_bytes: bytes) -> str:
    """Extrae texto de un archivo CSV."""
    text = _extract_from_txt(file_bytes)
    try:
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            rows.append(" | ".join(row))
        return "\n".join(rows) if rows else "[CSV vacío]"
    except Exception:
        return text
